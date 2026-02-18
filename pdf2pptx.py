import argparse
import os
import sys
from pathlib import Path
import tempfile
from PIL import Image
from PIL import ImageFilter
import fitz
from pptx import Presentation
from pptx.util import Inches


def clear_screen():
    """Clear the terminal screen based on the operating system."""
    os.system("cls")


def get_slide_dimensions(aspect_ratio, custom_width=None, custom_height=None) -> tuple[float, float]:
    """
    Get slide dimensions based on aspect ratio or custom dimensions.
    """
    if aspect_ratio == "custom" and custom_width and custom_height:
        return custom_width, custom_height

    # Standard aspect ratios (width, height) in inches
    ratios = {
        "4:3": (10.0, 7.5),  # Standard 4:3
        "16:9": (13.333, 7.5),  # Widescreen 16:9
        "16:10": (13.0, 8.125),  # 16:10 ratio
        "a4": (11.69, 8.27),  # A4 paper size
        "letter": (11.0, 8.5),  # Letter paper size
        "wide": (13.333, 7.5),  # Alias for 16:9
        "standard": (10.0, 7.5),  # Alias for 4:3
        "square": (10.0, 10.0),  # Square 1:1
        "portrait": (7.5, 10.0),  # Portrait orientation
    }

    return ratios.get(aspect_ratio, ratios["4:3"])


def convert_pdf_to_png(pdf_path, output_dir, page_num, dpi=300):
    """
    Convert PDF page to PNG format with specified DPI.

    Args:
        pdf_path: Path to PDF file
        output_dir: Output directory for PNG files
        page_num: Page number (0-indexed)
        dpi: Resolution for rendering

    Returns:
        Path to the created PNG file or None if failed
    """
    output_path = os.path.join(output_dir, f"slide-{page_num}.png")

    try:
        # Open PDF document
        pdf_doc = fitz.open(pdf_path)
        page = pdf_doc[page_num]

        # Render page as image
        pix = page.get_pixmap(dpi=dpi, alpha=False)

        # Save directly as PNG (no PIL conversion needed for basic saving)
        pix.save(output_path)

        pdf_doc.close()

        # Get image dimensions for logging
        # img_info = Image.open(output_path)
        # width, height = img_info.size
        # img_info.close()
        # print(f"  Page {page_num + 1}: {width}×{height} pixels ({dpi} DPI)")

        return output_path

    except Exception as e:
        print(f"Error converting page {page_num + 1} to PNG: {e}")
        return None


def convert_pdf_to_high_quality_png(pdf_path, output_dir, page_num, dpi=300, use_pil_enhancement=True):
    """
    Convert PDF page to high-quality PNG with optional PIL enhancements.

    Args:
        pdf_path: Path to PDF file
        output_dir: Output directory for PNG files
        page_num: Page number (0-indexed)
        dpi: Resolution for rendering
        use_pil_enhancement: Whether to use PIL for additional processing

    Returns:
        Path to the created PNG file or None if failed
    """
    output_path = os.path.join(output_dir, f"slide-{page_num}.png")

    try:
        # Open PDF document
        pdf_doc = fitz.open(pdf_path)
        page = pdf_doc[page_num]

        # Render page
        pix = page.get_pixmap(dpi=dpi, alpha=False)

        if use_pil_enhancement:
            # Convert to PIL Image for advanced processing
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

            # Optional: Apply sharpening for better text clarity
            img = img.filter(ImageFilter.SHARPEN)

            # Save with optimization
            img.save(output_path, "PNG",
                     dpi=(dpi, dpi),
                     optimize=True,
                     compress_level=6)
        else:
            # Save directly (faster)
            pix.save(output_path)

        pdf_doc.close()
        return output_path

    except Exception as e:
        print(f"Error converting page {page_num + 1} to PNG: {e}")
        return None

def pdf_to_pptx(pdf_path, output_path=None, aspect_ratio="4:3",
                custom_width=None, custom_height=None,
                quality="high", dpi=600, method="png"):
    """
    Convert PDF file to PPTX presentation without Ghostscript.

    Args:
        pdf_path: Path to the input PDF file
        output_path: Output PPTX file path (optional)
        aspect_ratio: Slide aspect ratio
        custom_width: Custom width in inches (if aspect_ratio is "custom")
        custom_height: Custom height in inches (if aspect_ratio is "custom")
        quality: Image quality ("low", "medium", "high", "ultra")
        dpi: Resolution for rendering (overrides quality setting if provided)
        method: Conversion method ("png")

    Returns:
        bool: True if conversion successful, False otherwise
    """
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"Error: File {pdf_path} not found.")
        return False

    # Output PPTX file (same name as PDF but with .pptx extension)
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pdf_path.with_suffix(".pptx")

    width_inches, height_inches = get_slide_dimensions(aspect_ratio, custom_width, custom_height)

    quality_settings = {
        "low": 150,
        "medium": 300,
        "high": 600,
        "ultra": 1200
    }

    if dpi is None and quality in quality_settings:
        dpi = quality_settings.get(quality, 600)

    # Create temporary directory for intermediate images
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir = Path(temp_dir)

        print(f"Processing: {pdf_path.name}")

        # Open PDF document
        pdf_doc = fitz.open(pdf_path)

        print(f"PDF pages: {len(pdf_doc)}")

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(width_inches)
        prs.slide_height = Inches(height_inches)
        # Use blank slide layout
        blank_slide_layout = prs.slide_layouts[6]

        # Process each PDF page
        for _, page in enumerate(pdf_doc.pages()):
            # clear_screen()
            # print(f"Processing page {_ + 1}/{len(pdf_doc)}...")

            sys.stdout.write("\rProcessing: Page %d/%d" % (_ + 1, len(pdf_doc)))
            sys.stdout.flush()

            # Render page as image with DPI
            pix = page.get_pixmap(dpi=dpi)

            # Convert to PIL Image
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

            # Resize image
            target_height = 1024
            ratio = target_height / img.height
            new_width = int(img.width * ratio)
            img = img.resize((new_width, target_height), Image.Resampling.LANCZOS)

            # Save image to temporary directory
            img_path = temp_dir / f"slide-{_}.png"
            img.save(img_path, "PNG")

            # Create slide and add image
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add image to fill entire slide
            left = top = Inches(0)
            slide.shapes.add_picture(str(img_path), left, top,
                                     width=prs.slide_width,
                                     height=prs.slide_height)

        # Close PDF document
        pdf_doc.close()

        # Save PowerPoint presentation
        prs.save(str(output_path))
        print("\r")
        print(f"Output file: {output_path}")
        print("Conversion completed!")

    return True

def pdf_to_pptx_simple(pdf_path, output_path=None, dpi=300):
    """
    Simple PDF to PPTX converter with minimal options.
    """
    return pdf_to_pptx(pdf_path=pdf_path, output_path=output_path, aspect_ratio="4:3", dpi=dpi, method="png")

def main():
    """Main function to handle command line arguments."""
    parser = argparse.ArgumentParser(
        description="Convert PDF files to PowerPoint presentations without Ghostscript."
    )

    parser.add_argument(
        "input_pdf",
        help="Path to the input PDF file"
    )

    parser.add_argument(
        "-o", "--output",
        help="Output PPTX file path (optional)"
    )

    parser.add_argument(
        "-a", "--aspect",
        choices=["4:3", "16:9", "16:10", "a4", "letter", "wide", "standard", "square", "portrait", "custom"],
        default="4:3",
        help="Slide aspect ratio (default: 4:3)"
    )

    parser.add_argument(
        "-w", "--width",
        type=float,
        help="Custom slide width in inches (requires --aspect custom)"
    )

    parser.add_argument(
        "-H", "--height",
        type=float,
        help="Custom slide height in inches (requires --aspect custom)"
    )

    parser.add_argument(
        "-q", "--quality",
        choices=["low", "medium", "high", "ultra"],
        default="high",
        help="Image quality (default: high)"
    )

    parser.add_argument(
        "-d", "--dpi",
        type=int,
        help="DPI for rendering (overrides --quality)"
    )

    parser.add_argument(
        "-m", "--method",
        choices=["png"],
        default="auto",
        help="Conversion method (default: auto), indev"
    )

    parser.add_argument(
        "--list-aspects",
        action="store_true",
        help="List available aspect ratios and exit"
    )

    parser.add_argument(
        "-s", "--simple",
        action="store_true",
        help="Use simple mode (4:3, 300 DPI)"
    )

    args = parser.parse_args()

    # List available aspect ratios if requested
    if args.list_aspects:
        print("Available aspect ratios:")
        print("  4:3      - Standard 4:3 (10.00 x 7.50 inches)")
        print("  16:9     - Widescreen 16:9 (13.33 x 7.50 inches)")
        print("  16:10    - 16:10 (13.00 x 8.125 inches)")
        print("  a4       - A4 paper size (11.69 x 8.27 inches)")
        print("  letter   - Letter paper size (11.00 x 8.50 inches)")
        print("  wide     - Alias for 16:9")
        print("  standard - Alias for 4:3")
        print("  square   - Square 1:1 (10.00 x 10.00 inches)")
        print("  portrait - Portrait orientation (7.50 x 10.00 inches)")
        print("  custom   - Custom dimensions (requires --width and --height)")
        sys.exit(0)

    # Validate custom dimensions
    if args.aspect == "custom":
        if args.width is None or args.height is None:
            print("Error: Custom aspect ratio requires both --width and --height arguments.")
            sys.exit(1)
        if args.width <= 0 or args.height <= 0:
            print("Error: Width and height must be positive numbers.")
            sys.exit(1)

    # Check input file
    pdf_path = Path(args.input_pdf)
    if not pdf_path.exists():
        print(f"Error: Input file '{pdf_path}' not found.")
        sys.exit(1)

    # Simple mode
    if args.simple:
        print("Running in simple mode (4:3, 300 DPI)...")
        success = pdf_to_pptx_simple(pdf_path, args.output, dpi=300)
    else:
        # Full mode
        success = pdf_to_pptx(
            pdf_path=pdf_path,
            output_path=args.output,
            aspect_ratio=args.aspect,
            custom_width=args.width,
            custom_height=args.height,
            quality=args.quality,
            dpi=args.dpi,
            method=args.method
        )

    # Exit with appropriate code
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()