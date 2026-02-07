import os
import platform
import sys
from pathlib import Path
import tempfile
from PIL import Image
import fitz
from pptx import Presentation
from pptx.util import Inches


def clear_screen():
    if platform.system() == "Windows":
        os.system('cls')
    else:
        os.system('clear')

def pdf_to_pptx(pdf_path):
    """
    Convert PDF file to PPTX presentation.

    Args:
        pdf_path: Path to the input PDF file

    Returns:
        bool: True if conversion successful, False otherwise
    """
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"Error: File {pdf_path} not found.")
        return False

    # Output PPTX file (same name as PDF but with .pptx extension)
    output_path = pdf_path.with_suffix('.pptx')

    # Create temporary directory for intermediate images
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir = Path(temp_dir)

        print(f"Processing: {pdf_path.name}")

        # Open PDF document
        pdf_doc = fitz.open(pdf_path)

        print(f"PDF pages: {len(pdf_doc)}")

        # Create PowerPoint presentation
        prs = Presentation()
        # Set to 4:3 widescreen aspect ratio
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        # Use blank slide layout
        blank_slide_layout = prs.slide_layouts[6]

        # Process each PDF page
        for i, page in enumerate(pdf_doc):
            # clear_screen()
            # print(f"Processing page {i + 1}/{len(pdf_doc)}...")
            sys.stdout.write("\rProcessing: Page %d/%d" % (i, len(pdf_doc)))
            sys.stdout.flush()

            # Render page as image with 300 DPI
            pix = page.get_pixmap(dpi=300)

            # Convert to PIL Image
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

            # Resize image
            target_height = 1024
            ratio = target_height / img.height
            new_width = int(img.width * ratio)
            img = img.resize((new_width, target_height), Image.Resampling.LANCZOS)

            # Save image to temporary directory
            img_path = temp_dir / f"slide-{i}.png"
            img.save(img_path, "PNG")

            # Create slide and add image
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add image to fill entire slide
            left = top = Inches(0)
            slide.shapes.add_picture(str(img_path), left, top,
                                     width=prs.slide_width,
                                     height=prs.slide_height)

        # Close PDF document
        pdf_doc.close()

        # Save PowerPoint presentation
        prs.save(str(output_path))
        print("\r")
        print(f"Output file: {output_path}")
        print("Conversion completed!")

    return True


if __name__ == "__main__":
    # Check command line arguments
    if len(sys.argv) < 2 or sys.argv[1] in ["help", "-h", "--help"]:
        print("Usage: python main.py filename.pdf")
        sys.exit(1)

    # Convert PDF to PPTX
    success = pdf_to_pptx(sys.argv[1])
    sys.exit(0 if success else 1)